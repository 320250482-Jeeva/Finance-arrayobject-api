"""
PPT Email Service API - Integrated with Existing Email Functions
Uses the user's existing code for authentication and email sending
"""
from fastapi import APIRouter, Request
import os
import json
import logging
from io import BytesIO
from typing import List, Dict, Any, Optional
from datetime import datetime
import mimetypes

import requests
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel, EmailStr, Field
from fastapi.responses import JSONResponse

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
router = APIRouter(tags=["Array Converter"])
# ============================================================================
# Logging
# ============================================================================

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# ============================================================================
# Configuration - From Your Existing Code
# ============================================================================

TOKEN_URL = "https://login.microsoftonline.com/1a407a2d-7675-4d17-8692-b3ac285306e4/oauth2/v2.0/token"
EMAIL_URL = "https://dev.apps.api.it.philips.com/api/email"

CLIENT_ID = os.getenv("CLIENT_ID", "826bc22b-bb13-471b-a9c3-10cfb0b11a83")
CLIENT_SECRET = os.getenv("CLIENT_SECRET", "Yy48Q~-5lglgh.GXF13sOp.qwYIFJq0-XPjPdc7O")
# os.getenv("CLIENT_SECRET") - if in .env load_dotenv()

SCOPE = os.getenv("SCOPE", "api://itaap-common-email-service-non-prod/.default")

MIME_OVERRIDES = {
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
}

# ============================================================================
# Pydantic Models
# ============================================================================

class PPTEmailRequest(BaseModel):
    """Request model for PPT generation and email"""
    business_name: str = Field(..., description="Business unit name", example="Philips EQ")
    summary: str = Field(..., description="Summary text", example="1. Finding\n2. Analysis")
    data: List[Dict[str, Any]] = Field(..., description="Data rows for table")
    email: EmailStr = Field(..., description="Recipient email")
    cc_emails: Optional[List[EmailStr]] = Field(default=None, description="CC emails")
    bcc_emails: Optional[List[EmailStr]] = Field(default=None, description="BCC emails")
    subject: Optional[str] = Field(default=None, description="Email subject")
    body: Optional[str] = Field(default=None, description="Email body HTML")


class APIResponse(BaseModel):
    """Response model"""
    success: bool
    message: str
    request_id: str
    timestamp: str
    pptx_filename: str
    email_status_code: Optional[int] = None


# ============================================================================
# PPT Mime Functions
# ============================================================================

def guess_mime(filename: str) -> str:
    """Get MIME type from filename"""
    ext = os.path.splitext(filename.lower())[1]
    if ext in MIME_OVERRIDES:
        return MIME_OVERRIDES[ext]
    mt, _ = mimetypes.guess_type(filename)
    return mt or "application/octet-stream"


def get_bearer_token(
    client_id: Optional[str] = None,
    client_secret: Optional[str] = None,
    scope: Optional[str] = None,
) -> str:
    """
    Your existing function - Generates OAuth2 bearer token
    From: AD Token generation (client_credentials flow)
    """
    client_id = client_id or CLIENT_ID
    client_secret = client_secret or CLIENT_SECRET
    scope = scope or SCOPE

    if not client_secret:
        raise RuntimeError("CLIENT_SECRET not set")

    data = {
        "grant_type": "client_credentials",
        "client_id": client_id,
        "scope": scope,
        "client_secret": client_secret,
    }

    resp = requests.post(TOKEN_URL, data=data, timeout=60)
    resp.raise_for_status()

    payload = resp.json()
    token = payload.get("access_token")
    if not token:
        raise RuntimeError(f"No access_token in response: {payload}")
    
    logger.info("Bearer token obtained")
    return token


def send_email(
    bearer_token: str,
    to_emails: List[str],
    cc_emails: Optional[List[str]] = None,
    bcc_emails: Optional[List[str]] = None,
    subject: str = "Test",
    body: str = "<html><body><h3>Test</h3></body></html>",
    attachment_buffer: Optional[BytesIO] = None,
    attachment_name: str = "attachment.pptx",
) -> requests.Response:
    """
    Your existing function - Sends email with attachment
    Uses Philips Email Service API
    """
    email_payload = {
        "to": to_emails,
        "cc": cc_emails or [],
        "bcc": bcc_emails or [],
        "subject": subject,
        "body": body,
        "platformName": "ITAAP",
        "projectName": "email-service",
        "priority": "1",
        "url": "email-attachments/",
    }

    headers = {
        "Authorization": f"Bearer {bearer_token}",
    }

    files = {
        "email-data": ("email-data.json", json.dumps(email_payload), "application/json"),
    }

    try:
        # Add attachment from memory
        if attachment_buffer:
            attachment_buffer.seek(0)
            files["attachment"] = (
                attachment_name,
                attachment_buffer,
                guess_mime(attachment_name)
            )

        resp = requests.post(EMAIL_URL, headers=headers, files=files, timeout=60)
        logger.info(f"Email sent with status: {resp.status_code}")
        return resp

    except requests.exceptions.RequestException as e:
        logger.error(f"Email sending failed: {e}")
        raise


# ============================================================================
# PPT Generation - v2.2 Integration
# ============================================================================

def create_pptx_buffer(business_name: str, summary: str, data: List[Dict[str, Any]]) -> BytesIO:
    """
    Create PPTX with v2.2 formatting:
    - Slide 1: Title (template background)
    - Slide 2: Summary (white background, centered)
    - Slide 3: Data table (white background, centered, negatives in red)
    - Slide 4: Thank you (template background)
    """
    logger.info(f"Generating PPTX for: {business_name}")
    
    prs = Presentation()
    
    # =====================================================================
    # SLIDE 1: Title (Template Background)
    # =====================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Remove placeholders
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)
    
    # Add title
    title_text = f"{business_name} - Analysis"
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(60)
    p.font.bold = True
    from pptx.dml.color import RGBColor
    p.font.color.rgb = RGBColor(79, 129, 189)  # Dark Blue Accent 1 Light 60%
    p.alignment = PP_ALIGN.CENTER
    
    # =====================================================================
    # SLIDE 2: Summary (White Background)
    # =====================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Set white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Remove placeholders
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)
    
    # Add Summary title
    title_left = Inches(0.75)
    title_top = Inches(0.5)
    title_width = Inches(8.5)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "SUMMARY"
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.LEFT
    
    # Parse and add summary
    bullets = _parse_summary(summary)
    
    content_left = Inches(0.75)
    content_top = Inches(1.5)
    content_width = Inches(7.5)
    content_height = Inches(4.5)
    
    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.clear()
    
    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = bullet_text
        p.font.size = Pt(14)
        p.level = 0
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.alignment = PP_ALIGN.LEFT
    
    # =====================================================================
    # SLIDE 3: Data Table (White Background)
    # =====================================================================
    if data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Set white background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Remove placeholders
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                sp = shape.element
                sp.getparent().remove(sp)
        
        # Add Data title
        title_left = Inches(0.75)
        title_top = Inches(0.5)
        title_width = Inches(8.5)
        title_height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = "DATA"
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.alignment = PP_ALIGN.LEFT
        
        # Create table
        keys = list(data[0].keys())
        num_rows = len(data) + 1
        num_cols = len(keys)
        
        table_width = Inches(8)
        left = Inches(0.5)
        top = Inches(1.5)
        height = Inches(4.5)
        
        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, height)
        table = table_shape.table
        
        # Set column widths
        col_width = table_width / num_cols
        for col in table.columns:
            col.width = int(col_width)
        
        # Header row
        for col_idx, key in enumerate(keys):
            cell = table.cell(0, col_idx)
            cell_text_frame = cell.text_frame
            cell_text_frame.clear()
            p = cell_text_frame.paragraphs[0]
            p.text = str(key)
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 51, 102)
        
        # Data rows
        for row_idx, row_data in enumerate(data, start=1):
            for col_idx, key in enumerate(keys):
                cell = table.cell(row_idx, col_idx)
                cell_value = str(row_data.get(key, ""))
                
                cell_text_frame = cell.text_frame
                cell_text_frame.clear()
                p = cell_text_frame.paragraphs[0]
                p.text = cell_value
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.CENTER
                
                # Red color for negative values
                if cell_value.strip().startswith('-'):
                    p.font.color.rgb = RGBColor(255, 0, 0)
                
                # Alternate row colors
                if row_idx % 2 == 0:
                    fill = cell.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(242, 242, 242)
    
    # =====================================================================
    # SLIDE 4: Thank You - Template Background
    # =====================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Remove placeholders
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)
    
    # Add thank you
    thank_you_left = Inches(0.5)
    thank_you_top = Inches(1.5)
    thank_you_width = Inches(8)
    thank_you_height = Inches(1.5)
    
    thank_you_box = slide.shapes.add_textbox(thank_you_left, thank_you_top, thank_you_width, thank_you_height)
    thank_you_frame = thank_you_box.text_frame
    thank_you_p = thank_you_frame.paragraphs[0]
    thank_you_p.text = "THANK YOU"
    thank_you_p.font.size = Pt(60)
    thank_you_p.font.bold = True
    thank_you_p.alignment = PP_ALIGN.CENTER
    
    # Save to BytesIO
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    
    logger.info("PPTX created successfully")
    return buf


def _parse_summary(summary: str) -> List[str]:
    """Parse summary into bullet points"""
    import re
    
    bullets = []
    clean = (summary or "").strip()
    
    if not clean:
        return ["No summary provided"]
    
    clean = re.sub(r"[ \t]+", " ", clean)
    clean = re.sub(r"\n\s*\n+", "\n", clean).strip()
    
    # Try numbered format
    numbered = re.findall(r"(?:^|\s)(\d+)\.\s", clean)
    if numbered:
        parts = re.split(r"(?:^|\s)(?=\d+\.\s)", clean)
        for part in parts:
            part = part.strip()
            if not part:
                continue
            part = re.sub(r"^\d+\.\s*", "", part).strip()
            if part:
                bullets.append(part)
    else:
        if "\n" in clean:
            bullets = [b.strip("-• ").strip() for b in clean.split("\n") if b.strip()]
        else:
            sentences = re.split(r"(?<=[.!?])\s+", clean)
            bullets = [s.strip() for s in sentences if s.strip()]
    
    return bullets if bullets else ["No summary provided"]


# ============================================================================
# FastAPI Application
# ============================================================================

@router.get("/health", tags=["Health"])
async def health_check():
    """Health check endpoint"""
    return {
        "status": "healthy",
        "service": "PPT Email Service API",
        "version": "1.0.0",
        "timestamp": datetime.now().isoformat()
    }


@router.post("/api/v1/generate-and-send", response_model=APIResponse, tags=["Email"])
async def generate_and_send(request: PPTEmailRequest) -> APIResponse:
    """
    Generate PPTX and send via email
    
    Parameters:
    - business_name: Business/company name
    - summary: Summary text (numbered format recommended)
    - data: List of dictionaries for table
    - email: Recipient email
    - cc_emails: CC recipients (optional)
    - bcc_emails: BCC recipients (optional)
    - subject: Email subject (optional, auto-generated if not provided)
    - body: Email body HTML (optional, auto-generated if not provided)
    """
    request_id = f"{datetime.now().strftime('%Y%m%d%H%M%S')}-{hash(request.email) % 10000}"
    timestamp = datetime.now().isoformat()
    
    logger.info(f"[{request_id}] Processing request for {request.business_name}")
    
    try:
        # Step 1: Generate PPTX
        logger.info(f"[{request_id}] Generating PPTX")
        pptx_buffer = create_pptx_buffer(
            business_name=request.business_name,
            summary=request.summary,
            data=request.data
        )
        
        pptx_filename = f"{request.business_name}_{request_id}.pptx"
        
        # Step 2: Get bearer token
        logger.info(f"[{request_id}] Authenticating")
        try:
            bearer_token = get_bearer_token()
        except Exception as e:
            logger.error(f"[{request_id}] Auth failed: {e}")
            raise HTTPException(status_code=401, detail=f"Authentication failed: {e}")
        
        # Step 3: Prepare email
        subject = request.subject or f"{request.business_name} - Analysis Report"
        body = request.body or f"""
        <html>
            <body>
                <h2>Report: {request.business_name}</h2>
                <p>Please find the analysis report attached.</p>
                <p><strong>Generated:</strong> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
            </body>
        </html>
        """
        
        # Step 4: Send email
        logger.info(f"[{request_id}] Sending email")
        try:
            email_response = send_email(
                bearer_token=bearer_token,
                to_emails=[request.email],
                cc_emails=request.cc_emails or [],
                bcc_emails=request.bcc_emails or [],
                subject=subject,
                body=body,
                attachment_buffer=pptx_buffer,
                attachment_name=pptx_filename
            )
            
            email_response.raise_for_status()
            
            logger.info(f"[{request_id}] Email sent successfully")
            
            return APIResponse(
                success=True,
                message="PPTX generated and email sent successfully",
                request_id=request_id,
                timestamp=timestamp,
                pptx_filename=pptx_filename,
                email_status_code=email_response.status_code
            )
            
        except requests.exceptions.HTTPError as e:
            logger.error(f"[{request_id}] Email error: {e}")
            raise HTTPException(
                status_code=e.response.status_code,
                detail=f"Email service error: {e.response.text}"
            )
        except Exception as e:
            logger.error(f"[{request_id}] Email sending failed: {e}")
            raise HTTPException(status_code=500, detail=f"Email sending failed: {e}")
            
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"[{request_id}] Unexpected error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/api/v1/example", tags=["Documentation"])
async def get_example():
    """Get example request/response"""
    return {
        "example_request": {
            "business_name": "Philips EQ",
            "summary": "1. Finding One\nAnalysis details.\n2. Finding Two\nMore analysis.",
            "data": [
                {"Market": "CEE", "YTD": "+45.2%", "Q4": "+28.5%"},
                {"Market": "RCA", "YTD": "-18.5%", "Q4": "-12.3%"},
            ],
            "email": "test@philips.com",
            "cc_emails": ["test@philips.com"],
            "subject": "Business Analysis - ITM 2026"
        },
        "example_response": {
            "success": True,
            "message": "PPTX generated and email sent successfully",
            "request_id": "20260112070551-1234",
            "timestamp": "2026-01-12T07:05:51.123456",
            "pptx_filename": "Philips EQ_20260112070551-1234.pptx",
            "email_status_code": 200
        }
    }
